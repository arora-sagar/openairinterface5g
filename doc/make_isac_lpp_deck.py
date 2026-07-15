#!/usr/bin/env python3
"""
Generate a plain (no-theme) Microsoft PowerPoint deck describing how ISAC
(Integrated Sensing and Communication) will be enabled on the OAI nrUE via
LPP, PRS and the T-tracer.

The content is grounded in the OpenAirInterface5G codebase:
  - openair3/LPP/MESSAGES/ASN1/37355-g60.asn        (LPP ASN.1, Rel-16, unused)
  - openair1/PHY/NR_UE_ESTIMATION/nr_dl_channel_estimation.c (PRS chest, ToA)
  - openair1/PHY/defs_nr_common.h                   (prs_meas_t: dl_toa, dl_aoa)
  - common/utils/T/T_messages.txt + T.h             (T-tracer events / T() macro)
  - radio/USRP/usrp_lib.cpp                          (GPSDO clock / 1-PPS sync)

Output: doc/ISAC_nrUE_LPP_Sensing.pptx  (plain white, no design theme)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Palette (neutral, print-friendly, no theme)
# ----------------------------------------------------------------------------
INK      = RGBColor(0x1A, 0x1A, 0x1A)   # near-black body text
ACCENT   = RGBColor(0x1F, 0x4E, 0x79)   # deep blue (headings / gNB)
SENSE    = RGBColor(0x8C, 0x3B, 0x00)   # burnt orange (sensing path)
COMMS    = RGBColor(0x1E, 0x6B, 0x3A)   # green (comms/positioning path)
UEBOX    = RGBColor(0x2E, 0x40, 0x57)   # slate (nrUE)
GREY     = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT    = RGBColor(0xF2, 0xF2, 0xF2)
EXISTS   = RGBColor(0x1E, 0x6B, 0x3A)   # green = exists today
TODO     = RGBColor(0x8C, 0x3B, 0x00)   # orange = to be added
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LINE     = RGBColor(0x7A, 0x7A, 0x7A)

FONT = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]        # fully blank layout -> no theme placeholders


# ----------------------------------------------------------------------------
# helpers
# ----------------------------------------------------------------------------
def _no_theme_bg(slide):
    """Force a solid white background so no template/theme fill shows."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE


def add_slide():
    s = prs.slides.add_slide(BLANK)
    _no_theme_bg(s)
    return s


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    return tb, tf


def set_run(r, text, size, color=INK, bold=False, italic=False, font=FONT):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color


def title_bar(slide, text, sub=None):
    """Heading + thin accent rule (no theme, just a rectangle rule)."""
    tb, tf = textbox(slide, Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.7))
    set_run(tf.paragraphs[0].add_run(), text, 26, ACCENT, bold=True)
    # accent rule
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(1.02), Inches(12.33), Pt(2.4))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    if sub:
        sb, sf = textbox(slide, Inches(0.5), Inches(1.06), Inches(12.3), Inches(0.4))
        set_run(sf.paragraphs[0].add_run(), sub, 12.5, GREY, italic=True)


def _indent(p, level):
    """Set left margin / hanging indent directly on the paragraph XML."""
    marL = int(Inches(0.24 + 0.28 * level))
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(marL))
    pPr.set('indent', str(-int(Inches(0.24))))


def bullet(tf, text, level=0, size=13.5, color=INK, bold=False, first=False,
           space_after=4, bullet_char="•"):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    _indent(p, level)
    # runs may be added by caller; here we add a single run
    r = p.add_run()
    set_run(r, ("%s  " % bullet_char if bullet_char else "") + text,
            size, color, bold=bold)
    return p


def rich_bullet(tf, segments, level=0, size=13.5, first=False, space_after=4,
                bullet_char="•", bullet_color=GREY):
    """segments: list of (text, color, bold)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    _indent(p, level)
    if bullet_char:
        r0 = p.add_run(); set_run(r0, bullet_char + "  ", size, bullet_color)
    for (txt, col, bold) in segments:
        r = p.add_run(); set_run(r, txt, size, col, bold=bold)
    return p


def box(slide, l, t, w, h, fill, text, tsize=12, tcolor=WHITE, bold=True,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_color=None, tanchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = tanchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); set_run(r, ln, tsize, tcolor, bold=(bold and i == 0))
        if i > 0:
            r.font.bold = False
    return sp


def connector(slide, x1, y1, x2, y2, color=LINE, width=1.75, arrow=True, dash=False):
    cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cxn.line.color.rgb = color
    cxn.line.width = Pt(width)
    ln = cxn.line._get_or_add_ln()
    if arrow:
        tail = ln.makeelement(qn('a:tailEnd'),
                              {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(d)
    return cxn


def small_label(slide, l, t, w, text, color=GREY, size=9.5, align=PP_ALIGN.CENTER,
                bold=False):
    tb, tf = textbox(slide, l, t, w, Inches(0.3))
    p = tf.paragraphs[0]; p.alignment = align
    set_run(p.add_run(), text, size, color, bold=bold)
    return tb


# ============================================================================
# SLIDE 1 — Title + context / motivation
# ============================================================================
s1 = add_slide()

# Title block
tb, tf = textbox(s1, Inches(0.6), Inches(1.15), Inches(12.1), Inches(1.5))
set_run(tf.paragraphs[0].add_run(),
        "Enabling ISAC on the OAI nrUE", 40, ACCENT, bold=True)
p = tf.add_paragraph()
set_run(p.add_run(),
        "Integrated Sensing & Communication via LPP, PRS and the T-tracer",
        20, UEBOX, bold=False)
p.space_before = Pt(6)

# Sub-context line
sb, sf = textbox(s1, Inches(0.6), Inches(2.55), Inches(12.1), Inches(0.5))
set_run(sf.paragraphs[0].add_run(),
        "Design overview — the nrUE as a joint sensing device and "
        "communication terminal", 14, GREY, italic=True)

# Two context columns
col_w = Inches(5.85)
# left: today
lb = box(s1, Inches(0.6), Inches(3.25), col_w, Inches(3.4), LIGHT, "",
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, tanchor=MSO_ANCHOR.TOP)
lb.line.color.rgb = LINE; lb.line.width = Pt(0.75)
_, ltf = textbox(s1, Inches(0.85), Inches(3.45), Inches(5.35), Inches(3.05))
rich_bullet(ltf, [("Today in OAI  ", ACCENT, True),
                  ("(baseline)", GREY, False)], size=15, first=True,
           bullet_char="")
bullet(ltf, "LPP (TS 37.355) exists as ASN.1 grammar only "
            "(openair3/LPP, Rel-16 g60) — compiled but not linked; no UE "
            "positioning entity.", size=12.5, color=INK)
bullet(ltf, "nrUE PRS reception computes per-antenna DL ToA, RSRP and SNR "
            "(nr_prs_channel_estimation).", size=12.5)
bullet(ltf, "CFO is estimated (PSS/PBCH) for compensation only; dl_aoa is a "
            "placeholder; per-antenna CIR is discarded.", size=12.5)
bullet(ltf, "Positioning stack (NRPPa / F1AP) lives on the gNB side, not the "
            "UE.", size=12.5)

# right: ISAC target
rb = box(s1, Inches(6.9), Inches(3.25), col_w, Inches(3.4), WHITE, "",
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, tanchor=MSO_ANCHOR.TOP)
rb.line.color.rgb = ACCENT; rb.line.width = Pt(1.25)
_, rtf = textbox(s1, Inches(7.15), Inches(3.45), Inches(5.35), Inches(3.05))
rich_bullet(rtf, [("ISAC target  ", ACCENT, True),
                  ("(this work)", GREY, False)], size=15, first=True,
           bullet_char="")
rich_bullet(rtf, [("UL/DL AoA, AoD & Doppler", SENSE, True),
                  (" derived at PHY and exported for sensing via the ",
                   INK, False),
                  ("T-tracer", SENSE, True), (".", INK, False)], size=12.5)
rich_bullet(rtf, [("PRS", SENSE, True),
                  (" reused as the sensing waveform — CIR → "
                   "range / angle / Doppler.", INK, False)], size=12.5)
rich_bullet(rtf, [("GPS & UE location / assistance", COMMS, True),
                  (" delivered over ", INK, False),
                  ("LPP", COMMS, True),
                  (" (A-GNSS, NR-DL-AoD, Multi-RTT, DL-TDOA).", INK, False)],
           size=12.5)
rich_bullet(rtf, [("One nrUE = ", INK, False),
                  ("communication + sensing", UEBOX, True),
                  (" on the same PRS/SRS signalling.", INK, False)], size=12.5)

# footer note about missing doc
fb, ff = textbox(s1, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4))
set_run(ff.paragraphs[0].add_run(),
        "Scope grounded in the OAI codebase and 3GPP Rel-16/17 ISAC & LPP "
        "concepts.", 10.5, GREY, italic=True)


# ============================================================================
# SLIDE 2 — Measurement & export pipeline (the "how")
# ============================================================================
s2 = add_slide()
title_bar(s2, "How AoA / AoD / Doppler & GPS are obtained",
          "PRS/SRS → PHY observables → T-tracer export (sensing)  |  "
          "measurements + GPS → LPP (communication / positioning)")

# Three method columns
top = Inches(1.35)
cw = Inches(4.02)
gap = Inches(0.14)
lefts = [Inches(0.5), Inches(0.5) + cw + gap, Inches(0.5) + 2 * (cw + gap)]
heads = ["1 . PRS → range & angle", "2 . Doppler / velocity",
         "3 . GPS time & UE location"]
hcols = [SENSE, SENSE, COMMS]

for l, head, hc in zip(lefts, heads, hcols):
    hb = box(s2, l, top, cw, Inches(0.52), hc, head, tsize=13.5, tcolor=WHITE)

# column 1 body
_, c1 = textbox(s2, lefts[0] + Inches(0.05), top + Inches(0.62),
                cw - Inches(0.1), Inches(3.15))
rich_bullet(c1, [("DL: ", ACCENT, True),
                 ("PRS channel estimate → CIR (IDFT). Peak = ", INK, False),
                 ("ToA / range", SENSE, True), (".", INK, False)],
           size=12, first=True)
rich_bullet(c1, [("Per-Rx-antenna complex CIR already computed → "
                  "inter-antenna phase gives ", INK, False),
                 ("DL-AoA / AoD", SENSE, True), (" (to add).", GREY, True)],
           size=12)
rich_bullet(c1, [("UL: ", ACCENT, True),
                 ("gNB SRS → ToA & UL-AoA (NRPPa / F1AP positioning).",
                  INK, False)], size=12)
rich_bullet(c1, [("Also available: ", INK, False),
                 ("RSRP, SNR", SENSE, True),
                 (" per antenna for detection gating.", INK, False)], size=12)

# column 2 body
_, c2 = textbox(s2, lefts[1] + Inches(0.05), top + Inches(0.62),
                cw - Inches(0.1), Inches(3.15))
rich_bullet(c2, [("CFO estimated at sync (", INK, False),
                 ("PSS fractional + PBCH fine", SENSE, True),
                 ("), stored as freq_offset.", INK, False)], size=12, first=True)
rich_bullet(c2, [("Phase drift of PRS CIR across occasions → ", INK, False),
                 ("Doppler shift", SENSE, True),
                 (" → radial velocity.", INK, False)], size=12)
rich_bullet(c2, [("freq↔velocity math already in NTN path "
                  "(dl_Doppler_shift).", INK, False)], size=12)
rich_bullet(c2, [("Coherent integration over PRS occasions → ", INK, False),
                 ("range–Doppler map", SENSE, True),
                 (" (to add).", GREY, True)], size=12)

# column 3 body
_, c3 = textbox(s2, lefts[2] + Inches(0.05), top + Inches(0.62),
                cw - Inches(0.1), Inches(3.15))
rich_bullet(c3, [("GPSDO on the RF front-end gives ", INK, False),
                 ("absolute GPS/TAI time + 1-PPS", COMMS, True),
                 (" (sync_to_gps).", INK, False)], size=12, first=True)
rich_bullet(c3, [("Timestamps every sensing measurement; aligns multi-gNB "
                  "PRS.", INK, False)], size=12)
rich_bullet(c3, [("UE geodetic location & assistance via ", INK, False),
                 ("LPP A-GNSS", COMMS, True),
                 (" (assistance data + navigation model).", INK, False)],
           size=12)
rich_bullet(c3, [("LPP ProvideLocationInformation carries NR-DL-AoD / "
                  "Multi-RTT / DL-TDOA + Velocity.", INK, False)], size=12)

# Export layer (bottom band): two rails
band_t = Inches(5.15)
# Sensing rail (T-tracer)
tb_sense = box(s2, Inches(0.5), band_t, Inches(6.02), Inches(1.75), WHITE, "",
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
tb_sense.line.color.rgb = SENSE; tb_sense.line.width = Pt(1.5)
_, sf2 = textbox(s2, Inches(0.7), band_t + Inches(0.1), Inches(5.65), Inches(1.55))
rich_bullet(sf2, [("SENSING export — T-tracer", SENSE, True)], size=13,
           first=True, bullet_char="")
rich_bullet(sf2, [("New event in T_messages.txt, e.g. ", INK, False),
                  ("T_UE_PHY_SENSING", SENSE, True),
                  (" : frame:slot:aoa:aod:doppler:toa:rsrp.", INK, False)],
           size=11.5)
rich_bullet(sf2, [("Emitted with the ", INK, False),
                  ("T(…)", SENSE, True),
                  (" macro at the PRS chest site (as existing "
                   "T_UE_PHY_DL_CHANNEL_ESTIMATE).", INK, False)], size=11.5)
rich_bullet(sf2, [("Collected live by ", INK, False),
                  ("record / textlog / extract", SENSE, True),
                  (" → offline range–Doppler / angle processing.",
                   INK, False)], size=11.5)

# Comms/positioning rail (LPP)
tb_comm = box(s2, Inches(6.81), band_t, Inches(6.02), Inches(1.75), WHITE, "",
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
tb_comm.line.color.rgb = COMMS; tb_comm.line.width = Pt(1.5)
_, cf2 = textbox(s2, Inches(7.01), band_t + Inches(0.1), Inches(5.65), Inches(1.55))
rich_bullet(cf2, [("COMMUNICATION / positioning export — LPP", COMMS, True)],
           size=13, first=True, bullet_char="")
rich_bullet(cf2, [("New UE LPP entity links the existing ", INK, False),
                  ("asn1_lpp", COMMS, True),
                  (" codec (37355-g60).", INK, False)], size=11.5)
rich_bullet(cf2, [("Transports LPP transparently over NAS/RRC to the ", INK, False),
                  ("LMF", COMMS, True), (".", INK, False)], size=11.5)
rich_bullet(cf2, [("Capabilities → Assistance → "
                   "ProvideLocationInformation with the measured "
                   "AoA/AoD/ToA/Doppler + GNSS.", INK, False)], size=11.5)


# ============================================================================
# SLIDE 3 — Architecture diagram: nrUE as sensing + communication node
# ============================================================================
s3 = add_slide()
title_bar(s3, "nrUE as a joint sensing + communication node",
          "One PHY, two outputs: sensing observables (T-tracer) and "
          "positioning/comms (LPP to LMF)")

# --- gNB / TRP (left) ---
gnb = box(s3, Inches(0.45), Inches(2.55), Inches(1.85), Inches(1.5), ACCENT,
          "gNB / TRP", tsize=14)
small_label(s3, Inches(0.30), Inches(4.10), Inches(2.15),
            "DL PRS  /  UL SRS", ACCENT, 10.5, bold=True)
# target/object for sensing (echo)
tgt = box(s3, Inches(0.62), Inches(1.55), Inches(1.6), Inches(0.72), GREY,
          "Object /\ntarget", tsize=11, shape=MSO_SHAPE.OVAL)

# --- nrUE PHY (center) ---
ue = box(s3, Inches(3.05), Inches(2.25), Inches(3.15), Inches(2.15), UEBOX, "",
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
_, uef = textbox(s3, Inches(3.2), Inches(2.35), Inches(2.9), Inches(2.0),
                 anchor=MSO_ANCHOR.TOP)
p = uef.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "OAI nrUE  —  PHY", 14, WHITE, bold=True)
for txt in ["PRS channel est. → CIR",
            "ToA / RSRP / SNR",
            "per-antenna CIR (AoA/AoD)",
            "CFO → Doppler"]:
    pp = uef.add_paragraph(); pp.alignment = PP_ALIGN.CENTER
    pp.space_after = Pt(1)
    set_run(pp.add_run(), txt, 11, WHITE, bold=False)
# GPSDO chip on UE
gps = box(s3, Inches(3.05), Inches(4.55), Inches(3.15), Inches(0.5), COMMS,
          "GPSDO: GPS/TAI time + 1-PPS", tsize=10.5, tcolor=WHITE)

# --- Sensing output (top right) ---
ttr = box(s3, Inches(7.35), Inches(1.35), Inches(2.7), Inches(1.15), WHITE,
          "T-tracer\nT_UE_PHY_SENSING", tsize=12.5, tcolor=SENSE)
ttr.line.color.rgb = SENSE; ttr.line.width = Pt(1.5)
sens = box(s3, Inches(10.4), Inches(1.35), Inches(2.5), Inches(1.15), LIGHT,
           "SENSING\nrange · AoA/AoD · Doppler\nrange–Doppler map",
           tsize=11, tcolor=SENSE)
sens.line.color.rgb = SENSE; sens.line.width = Pt(1)

# --- Comms/positioning output (bottom right) ---
lpp = box(s3, Inches(7.35), Inches(3.85), Inches(2.7), Inches(1.15), WHITE,
          "LPP entity\n(asn1_lpp / TS 37.355)", tsize=12, tcolor=COMMS)
lpp.line.color.rgb = COMMS; lpp.line.width = Pt(1.5)
lmf = box(s3, Inches(10.4), Inches(3.85), Inches(2.5), Inches(1.15), LIGHT,
          "LMF\nposition fix +\ncommunication",
          tsize=11.5, tcolor=COMMS)
lmf.line.color.rgb = COMMS; lmf.line.width = Pt(1)
small_label(s3, Inches(7.35), Inches(5.02), Inches(2.7),
            "transparent over NAS / RRC", COMMS, 9.5)

# --- connectors ---
E = Emu
# gNB <-> UE (DL PRS / UL SRS) bidirectional-ish (two arrows)
connector(s3, E(Inches(2.30)), E(Inches(3.05)), E(Inches(3.05)), E(Inches(3.05)),
          color=ACCENT, width=2.0, arrow=True)
connector(s3, E(Inches(3.05)), E(Inches(3.55)), E(Inches(2.30)), E(Inches(3.55)),
          color=ACCENT, width=2.0, arrow=True, dash=True)
# object echo path (target reflects PRS back to UE) -> sensing
connector(s3, E(Inches(1.75)), E(Inches(2.27)), E(Inches(3.10)), E(Inches(2.70)),
          color=SENSE, width=1.75, arrow=True, dash=True)
small_label(s3, Inches(2.28), Inches(1.62), Inches(1.55),
            "echo", SENSE, 9, align=PP_ALIGN.LEFT)
# UE -> T-tracer (sensing)
connector(s3, E(Inches(6.20)), E(Inches(2.9)), E(Inches(7.35)), E(Inches(1.95)),
          color=SENSE, width=2.0, arrow=True)
# T-tracer -> sensing outputs
connector(s3, E(Inches(10.05)), E(Inches(1.92)), E(Inches(10.4)), E(Inches(1.92)),
          color=SENSE, width=2.0, arrow=True)
# UE -> LPP (comms/positioning)
connector(s3, E(Inches(6.20)), E(Inches(3.7)), E(Inches(7.35)), E(Inches(4.42)),
          color=COMMS, width=2.0, arrow=True)
# LPP -> LMF
connector(s3, E(Inches(10.05)), E(Inches(4.42)), E(Inches(10.4)), E(Inches(4.42)),
          color=COMMS, width=2.0, arrow=True)
# GPSDO -> LPP (time/location) and -> T-tracer (timestamp) small dashed
connector(s3, E(Inches(4.6)), E(Inches(5.05)), E(Inches(4.6)), E(Inches(5.05)),
          color=COMMS, width=1.0, arrow=False)

# path labels
small_label(s3, Inches(6.15), Inches(2.25), Inches(1.4), "sensing", SENSE, 10,
            bold=True)
small_label(s3, Inches(6.15), Inches(3.95), Inches(1.4), "positioning",
            COMMS, 10, bold=True)

# --- legend / status table (bottom) ---
leg_t = Inches(5.55)
_, lf = textbox(s3, Inches(0.45), leg_t, Inches(6.6), Inches(1.6))
rich_bullet(lf, [("Status", ACCENT, True)], size=12.5, first=True, bullet_char="")
rich_bullet(lf, [("Exists today: ", EXISTS, True),
                 ("PRS ToA/RSRP/SNR, per-antenna CIR, CFO, LPP ASN.1, GPSDO "
                  "time, T-tracer.", INK, False)], size=11.5,
           bullet_char="✓", bullet_color=EXISTS)
rich_bullet(lf, [("To add for ISAC: ", TODO, True),
                 ("AoA/AoD & Doppler estimators, sensing T-event, UE LPP "
                  "entity + NAS/RRC transport.", INK, False)], size=11.5,
           bullet_char="▸", bullet_color=TODO)

_, lf2 = textbox(s3, Inches(7.2), leg_t, Inches(5.6), Inches(1.6))
rich_bullet(lf2, [("Legend", ACCENT, True)], size=12.5, first=True,
           bullet_char="")
rich_bullet(lf2, [("sensing data path — T-tracer "
                   "(range, AoA/AoD, Doppler).", INK, False)],
           size=11.5, bullet_char="■", bullet_color=SENSE)
rich_bullet(lf2, [("communication / positioning path — LPP, GPS.",
                   INK, False)],
           size=11.5, bullet_char="■", bullet_color=COMMS)
rich_bullet(lf2, [("PRS (DL) / SRS (UL) radio link.", INK, False)],
           size=11.5, bullet_char="■", bullet_color=ACCENT)

# ----------------------------------------------------------------------------
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "ISAC_nrUE_LPP_Sensing.pptx")
prs.save(out)
print("Saved:", out, "slides:", len(prs.slides._sldIdLst))
